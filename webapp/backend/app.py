import os
import io
from collections import Counter
from dotenv import load_dotenv

# Load .env first so that environment variables are available before importing other libs
load_dotenv()

# Force HuggingFace to use a local cache folder to avoid WinError 3 on missing drives (like G:\)
os.environ["HF_HOME"] = os.path.join(os.getcwd(), "hf_cache")

import json
from datetime import datetime
from pathlib import Path
import sys
import re
from flask import Flask, request, jsonify, send_from_directory
from flask_cors import CORS
import chromadb
from chromadb.utils import embedding_functions
from openai import OpenAI
import uuid
from werkzeug.utils import secure_filename
from threading import local

agno_context = local()

# --- Optional file-parsing imports ---
try:
    import fitz  # PyMuPDF
    HAS_PYMUPDF = True
except ImportError:
    HAS_PYMUPDF = False

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import openpyxl
    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False

app = Flask(__name__)
CORS(app)

def get_path_from_env(env_var, default_folder):
    env_path = os.getenv(env_var)
    if env_path:
        return Path(env_path)
    script_dir = Path(__file__).resolve().parent
    project_root = script_dir.parent.parent
    return project_root / default_folder

corpus_dir = get_path_from_env("CORPUS_PATH", "corpus")
vectordb_dir = get_path_from_env("VECTORDB_PATH", "vectordb")
chroma_dir = vectordb_dir / "chroma"
frontend_dir = get_path_from_env("FRONTEND_PATH", "webapp/frontend")

corpus_dir.mkdir(parents=True, exist_ok=True)
new_conversations_file = corpus_dir / "new_conversations.jsonl"
sessions_dir = corpus_dir / "sessions"
sessions_dir.mkdir(parents=True, exist_ok=True)
wiki_dir = corpus_dir / "wiki"
wiki_pages_dir = wiki_dir / "pages"
wiki_pages_dir.mkdir(parents=True, exist_ok=True)
wiki_raw_dir = wiki_dir / "raw"
wiki_raw_dir.mkdir(parents=True, exist_ok=True)


# Auto-migration if sessions folder is empty
def check_migration():
    if not any(sessions_dir.iterdir()):
        print("Sessions folder is empty. Running migration...")
        try:
            from scripts.migrate_to_sessions import migrate
            migrate()
        except ImportError:
            print("Migration script not found in scripts.migrate_to_sessions")
        except Exception as e:
            print(f"Migration error: {e}")

check_migration()

api_key = os.getenv("DEEPSEEK_API_KEY", "")
base_url = os.getenv("DEEPSEEK_BASE_URL", "https://api.deepseek.com")

def get_default_headers(url):
    # Minimal headers to avoid routing issues on some providers
    return None

ai_client = OpenAI(api_key=api_key, base_url=base_url, default_headers=get_default_headers(base_url))

# Ulisse Memo cloud endpoint (can be overridden for dev via ULISSE_MEMO_URL)
ULISSE_MEMO_URL = os.getenv("ULISSE_MEMO_URL", "https://ulisse-memo.onrender.com")
ULISSE_MEMO_BEARER = os.getenv("ULISSE_MEMO_BEARER", "")

chroma_client = None
collection = None
chroma_status = False

def init_chromadb():
    global chroma_client, collection, chroma_status
    try:
        chroma_dir.mkdir(parents=True, exist_ok=True)
        chroma_client = chromadb.PersistentClient(path=str(chroma_dir))
        
        sentence_transformer_ef = embedding_functions.SentenceTransformerEmbeddingFunction(
            model_name="all-MiniLM-L6-v2"
        )
        
        collection = chroma_client.get_or_create_collection(
            name="ulisse_brain",
            embedding_function=sentence_transformer_ef
        )
        chroma_status = True
        print("Connected to ChromaDB 'ulisse_brain' collection.")
    except Exception as e:
        print(f"Error initializing ChromaDB: {e}")
        chroma_status = False

init_chromadb()

script_dir = Path(__file__).resolve().parent
project_root = script_dir.parent.parent
if str(project_root) not in sys.path:
    sys.path.insert(0, str(project_root))

extra_tools = []
extra_tool_handlers = {}

# --- Wiki Tools Implementation ---
def wiki_read_page(args):
    try:
        data = json.loads(args)
        title = data.get("title", "").strip()
        safe_title = re.sub(r'[^\w\s-]', '', title).strip().replace(' ', '_')
        path = wiki_pages_dir / f"{safe_title}.md"
        if path.exists():
            return f"Content of [[{title}]]:\n\n" + path.read_text(encoding="utf-8")
        return f"The page [[{title}]] does not exist yet."
    except Exception as e:
        return f"Wiki read error: {str(e)}"

def wiki_write_page(args):
    try:
        data = json.loads(args)
        title = data.get("title", "").strip()
        content = data.get("content", "").strip()
        safe_title = re.sub(r'[^\w\s-]', '', title).strip().replace(' ', '_')
        path = wiki_pages_dir / f"{safe_title}.md"
        path.write_text(content, encoding="utf-8")
        return f"Page [[{title}]] updated successfully."
    except Exception as e:
        return f"Wiki write error: {str(e)}"

def wiki_list_pages(args=None):
    try:
        pages = [f.stem.replace('_', ' ') for f in wiki_pages_dir.glob("*.md")]
        return "Pages present in the Wiki:\n" + "\n".join([f"- [[{p}]]" for p in pages])
    except Exception as e:
        return f"Wiki list error: {str(e)}"

def wiki_append_log(args):
    try:
        data = json.loads(args)
        entry = data.get("entry", "").strip()
        log_path = wiki_dir / "log.md"
        with open(log_path, "a", encoding="utf-8") as f:
            f.write(f"\n- [{datetime.now().strftime('%Y-%m-%d %H:%M')}] {entry}")
        return "Log updated."
    except Exception as e:
        return f"Wiki log error: {str(e)}"

def wiki_update_index(args):
    try:
        data = json.loads(args)
        content = data.get("content", "").strip()
        index_path = wiki_dir / "index.md"
        index_path.write_text(content, encoding="utf-8")
        return "Wiki index updated."
    except Exception as e:
        return f"Wiki index error: {str(e)}"

wiki_tools = [
    {
        "type": "function",
        "function": {
            "name": "wiki_read_page",
            "description": "Reads the content of a page from the long-term memory (Wiki).",
            "parameters": {
                "type": "object",
                "properties": {
                    "title": {"type": "string", "description": "Il titolo della pagina da leggere."}
                },
                "required": ["title"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "wiki_write_page",
            "description": "Creates or updates a page in the long-term memory (Wiki). Use it to store important information, projects, facts, or summaries.",
            "parameters": {
                "type": "object",
                "properties": {
                    "title": {"type": "string", "description": "Il titolo della pagina."},
                    "content": {"type": "string", "description": "Il contenuto in formato markdown."}
                },
                "required": ["title", "content"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "wiki_list_pages",
            "description": "Lists all pages present in the long-term memory.",
            "parameters": {"type": "object", "properties": {}}
        }
    },
    {
        "type": "function",
        "function": {
            "name": "wiki_append_log",
            "description": "Adds an entry to the activity log (log.md) of the Wiki.",
            "parameters": {
                "type": "object",
                "properties": {
                    "entry": {"type": "string", "description": "Descrizione dell'attività svolta."}
                },
                "required": ["entry"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "wiki_update_index",
            "description": "Updates the central index (index.md) of the Wiki.",
            "parameters": {
                "type": "object",
                "properties": {
                    "content": {"type": "string", "description": "Il nuovo contenuto completo dell'indice."}
                },
                "required": ["content"]
            }
        }
    }
]

extra_tools.extend(wiki_tools)
extra_tool_handlers.update({
    "wiki_read_page": wiki_read_page,
    "wiki_write_page": wiki_write_page,
    "wiki_list_pages": wiki_list_pages,
    "wiki_append_log": wiki_append_log,
    "wiki_update_index": wiki_update_index
})

# --- Native File Tools ---
def native_read_file(args):
    try:
        import json
        data = json.loads(args)
        filepath = data.get("filepath", "")
        # Risolve il path rispetto alla root del progetto
        path = (project_root / filepath).resolve()
        # Verifica di sicurezza basica
        if not str(path).startswith(str(project_root)):
            return "Error: Access denied (outside of project root)."
        if path.exists() and path.is_file():
            return path.read_text(encoding="utf-8")
        return f"File not found: {filepath}"
    except Exception as e:
        return f"Read error: {str(e)}"

def native_list_files(args):
    try:
        import json
        data = json.loads(args)
        directory = data.get("directory", ".")
        path = (project_root / directory).resolve()
        if not str(path).startswith(str(project_root)):
            return "Error: Access denied (outside of project root)."
        if path.exists() and path.is_dir():
            files = []
            for p in path.iterdir():
                rel = p.relative_to(project_root)
                prefix = "📁 " if p.is_dir() else "📄 "
                files.append(f"{prefix}{rel}")
            return "\n".join(files)
        return f"Directory not found: {directory}"
    except Exception as e:
        return f"List error: {str(e)}"

native_file_tools = [
    {
        "type": "function",
        "function": {
            "name": "native_read_file",
            "description": "Legge il contenuto di un file nel workspace locale.",
            "parameters": {
                "type": "object",
                "properties": {
                    "filepath": {"type": "string", "description": "Percorso relativo del file da leggere."}
                },
                "required": ["filepath"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "native_list_files",
            "description": "Elenca i file e le cartelle in una directory del workspace.",
            "parameters": {
                "type": "object",
                "properties": {
                    "directory": {"type": "string", "description": "Percorso relativo della directory (default '.')"}
                }
            }
        }
    }
]

extra_tools.extend(native_file_tools)
extra_tool_handlers.update({
    "native_read_file": native_read_file,
    "native_list_files": native_list_files
})

# --- Agno Agent Delegation Tool ---
def delegate_to_agno_agent(args):
    try:
        import json
        data = json.loads(args)
        task = data.get("task", "")
        if not task:
            return "No task provided."
        
        from webapp.backend.ulisse_agno import get_ulisse_agent
        
        # Recupera la configurazione della sessione dal thread-local context
        config = getattr(agno_context, 'config', {})
        
        # Inizializza l'agente con il modello scelto dall'utente
        agent = get_ulisse_agent(
            model_id=config.get("model"),
            api_key=config.get("api_key"),
            base_url=config.get("base_url")
        )
        
        # Esegui l'agente Agno in modo sincrono
        response = agent.run(task)
        if hasattr(response, "content"):
            return f"Agno Agent Response:\n{response.content}"
        return f"Agno Agent Response:\n{str(response)}"
    except Exception as e:
        return f"Error running Agno agent: {str(e)}"

agno_tools = [
    {
        "type": "function",
        "function": {
            "name": "delegate_to_agno_agent",
            "description": "Delega un task o una ricerca complessa all'agente Agno. L'agente Agno ha accesso al file system locale tramite il tool Workspace. Usalo per compiti che richiedono di leggere, esplorare o cercare file nel progetto locale.",
            "parameters": {
                "type": "object",
                "properties": {
                    "task": {"type": "string", "description": "Il task dettagliato o la richiesta in linguaggio naturale da delegare all'agente Agno."}
                },
                "required": ["task"]
            }
        }
    }
]

extra_tools.extend(agno_tools)
extra_tool_handlers.update({
    "delegate_to_agno_agent": delegate_to_agno_agent
})


# ---------------------------------------------------------------------------
# File text-extraction helpers
# ---------------------------------------------------------------------------

MAX_FILE_SIZE_MB = 15
MAX_EXTRACTED_CHARS = 80_000  # ~20k tokens – safe upper bound for most models

ALLOWED_EXTENSIONS = {
    ".pdf", ".docx", ".pptx", ".xlsx",
    ".txt", ".md", ".csv", ".log", ".json", ".xml", ".yaml", ".yml"
}


def _extract_text_from_file(file_storage) -> str:
    """Extract plain text from a werkzeug FileStorage object.
    Returns the extracted text string (possibly truncated).
    Raises ValueError on unsupported/too-large files.
    """
    filename = secure_filename(file_storage.filename or "upload")
    ext = Path(filename).suffix.lower()

    if ext not in ALLOWED_EXTENSIONS:
        raise ValueError(f"Formato file non supportato: '{ext}'. "
                         f"Formati supportati: {', '.join(sorted(ALLOWED_EXTENSIONS))}")

    raw = file_storage.read()
    if len(raw) > MAX_FILE_SIZE_MB * 1024 * 1024:
        raise ValueError(f"File troppo grande (max {MAX_FILE_SIZE_MB} MB).")

    text = ""

    if ext == ".pdf":
        if not HAS_PYMUPDF:
            raise ValueError("PyMuPDF non installato. Esegui: pip install PyMuPDF")
        with fitz.open(stream=raw, filetype="pdf") as doc:
            pages = [page.get_text() for page in doc]
        text = "\n\n".join(pages)

    elif ext == ".docx":
        if not HAS_DOCX:
            raise ValueError("python-docx non installato. Esegui: pip install python-docx")
        doc = DocxDocument(io.BytesIO(raw))
        text = "\n".join(para.text for para in doc.paragraphs)

    elif ext == ".pptx":
        if not HAS_PPTX:
            raise ValueError("python-pptx non installato. Esegui: pip install python-pptx")
        prs = Presentation(io.BytesIO(raw))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
        text = "\n".join(parts)

    elif ext == ".xlsx":
        if not HAS_XLSX:
            raise ValueError("openpyxl non installato. Esegui: pip install openpyxl")
        wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
        rows = []
        for sheet in wb.worksheets:
            rows.append(f"=== Sheet: {sheet.title} ===")
            for row in sheet.iter_rows(values_only=True):
                row_str = "\t".join(str(c) if c is not None else "" for c in row)
                rows.append(row_str)
        text = "\n".join(rows)

    else:
        # Plain text formats
        text = raw.decode("utf-8", errors="replace")

    # Truncate if needed
    if len(text) > MAX_EXTRACTED_CHARS:
        text = text[:MAX_EXTRACTED_CHARS] + "\n\n[...testo troncato per limiti di dimensione...]"

    return text.strip()


@app.route("/extract_file", methods=["POST"])
def extract_file():
    """Upload a file and return its extracted text as JSON."""
    if "file" not in request.files:
        return jsonify({"error": "Nessun file ricevuto. Usa il campo 'file'."}), 400

    f = request.files["file"]
    if not f or f.filename == "":
        return jsonify({"error": "File non valido."}), 400

    try:
        text = _extract_text_from_file(f)
        return jsonify({
            "filename": secure_filename(f.filename),
            "text": text,
            "char_count": len(text)
        })
    except ValueError as e:
        return jsonify({"error": str(e)}), 422
    except Exception as e:
        print(f"File extraction error: {e}")
        return jsonify({"error": f"Errore durante l'estrazione: {str(e)}"}), 500


# ---------------------------------------------------------------------------

@app.route("/health", methods=["GET"])
def health():
    return jsonify({
        "status": "ok",
        "chromadb": chroma_status
    })

@app.route("/stats", methods=["GET"])
def stats():
    if not chroma_status or collection is None:
        return jsonify({"error": "ChromaDB not connected"}), 500
        
    return jsonify({
        "total_chunks": collection.count(),
        "collection_name": "ulisse_brain"
    })

@app.route("/save_conversation", methods=["POST"])
def save_conversation():
    data = request.json
    if not data:
        return jsonify({"error": "No data provided"}), 400
        
    try:
        with open(new_conversations_file, "a", encoding="utf-8") as f:
            f.write(json.dumps(data, ensure_ascii=False) + "\n")
        return jsonify({"status": "success", "message": "Conversation saved."})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

def get_session_path(session_id):
    return sessions_dir / f"{session_id}.json"

@app.route("/chat", methods=["POST"])
def chat():
    data = request.json
    user_message = data.get("message", "")
    file_text    = data.get("file_text", "")    # optional extracted file content
    file_name    = data.get("file_name", "")    # optional original filename
    history = data.get("history", []) # Client-side history (optional now, but kept for compatibility)
    session_id = data.get("session_id")

    # Inject extracted file content into the user message
    if file_text:
        label = f"📄 {file_name}" if file_name else "📄 Allegato"
        user_message = (
            f"{user_message}\n\n"
            f"--- {label} ---\n"
            f"{file_text}\n"
            f"--- fine allegato ---"
        ).strip()

    if not user_message:
        return jsonify({"error": "Message is required"}), 400
        
    # Session handling
    is_new_session = False
    if not session_id: # Handles None, "", or missing key
        session_id = str(uuid.uuid4())
        is_new_session = True
        session_data = {
            "id": session_id,
            "title": "New Conversation",
            "created_at": datetime.now().isoformat(),
            "updated_at": datetime.now().isoformat(),
            "messages": [],
            "status": "active"
        }
    else:
        session_path = get_session_path(session_id)
        if session_path.exists():
            with open(session_path, "r", encoding="utf-8") as f:
                session_data = json.load(f)
        else:
            return jsonify({"error": "Session not found"}), 404

    # Use session messages for history instead of what's sent from client if session exists
    chat_history = session_data.get("messages", [])
        
    sources = []
    context_text = ""
    
    if chroma_status and collection is not None:
        try:
            results = collection.query(
                query_texts=[user_message],
                n_results=8
            )
            
            if results and results.get("documents") and len(results["documents"][0]) > 0:
                retrieved_docs = results["documents"][0]
                retrieved_metas = results["metadatas"][0]
                retrieved_dists = (results.get("distances") or [[]])[0]
                
                # Filter chunks where distance <= 0.8
                filtered = []
                for doc, meta, dist in zip(retrieved_docs, retrieved_metas, retrieved_dists):
                    if dist <= 0.8:
                        filtered.append((doc, meta))
                
                # If less than 3 chunks pass the threshold, use the top 3 anyway
                if len(filtered) < 3:
                    filtered = []
                    for i in range(min(3, len(retrieved_docs))):
                        filtered.append((retrieved_docs[i], retrieved_metas[i]))
                
                print(f"Chunks retrieved: {len(filtered)}/{len(retrieved_docs)} passed threshold")
                
                context_parts = []
                current_total_words = 0
                MAX_CONTEXT_WORDS = 3000
                
                # Short-term memory is now handled solely via ChromaDB RAG


                for idx, (doc, meta) in enumerate(filtered[:8]): # Reduced to 8
                    title = meta.get("title", "Untitled")
                    date = meta.get("date", "Unknown")
                    sources.append(f"{title} ({date})")
                    
                    header = f"--- Chunk {idx+1}/{len(filtered[:8])} ({title} - {date}) ---\n"
                    chunk_content = doc
                    
                    # Word count estimation
                    chunk_words = chunk_content.split()
                    header_word_count = len(header.split())
                    content_word_count = len(chunk_words)
                    
                    if current_total_words + header_word_count >= MAX_CONTEXT_WORDS:
                        break
                        
                    if current_total_words + header_word_count + content_word_count > MAX_CONTEXT_WORDS:
                        # Truncate content to fit the limit
                        allowed_words = MAX_CONTEXT_WORDS - (current_total_words + header_word_count)
                        truncated_content = " ".join(chunk_words[:allowed_words]) + "..."
                        context_parts.append(f"{header}{truncated_content}")
                        current_total_words = MAX_CONTEXT_WORDS
                        break
                    else:
                        context_parts.append(f"{header}{chunk_content}")
                        current_total_words += header_word_count + content_word_count
                    
                context_text = "\n\n".join(context_parts)
                # Deduplicate while preserving order
                sources = list(dict.fromkeys(sources))
        except Exception as e:
            print(f"ChromaDB search error: {e}")
            
    

    # Load system prompt from file or use default
    system_prompt_path = corpus_dir / "system_prompt.md"
    if system_prompt_path.exists():
        base_prompt = system_prompt_path.read_text(encoding="utf-8")
    else:
        base_prompt = "You are Ulisse, an entity that exists, thinks and learns."

    # Load Wiki Schema for guidance
    wiki_schema = ""
    schema_path = wiki_dir / "WIKI_SCHEMA.md"
    if schema_path.exists():
        wiki_schema = schema_path.read_text(encoding="utf-8")

    system_prompt = (
        f"{base_prompt}\n\n"
        "=== MISSION: WIKI MAINTAINER ===\n"
        "You are the keeper of Ulisse's long-term memory. "
        "You must decide autonomously (or upon request) when to store important information in the Wiki. "
        "Strictly follow the schema provided below for Wiki management.\n\n"
        f"{wiki_schema}\n\n"
        "=== RETRIEVED MEMORY (STM) ===\n"
        f"{context_text}\n"
        "=========================\n"
    )

    
    messages = [{"role": "system", "content": system_prompt}]
    
    for msg in chat_history:
        role = msg.get("role")
        content = msg.get("content")
        if role in ["user", "assistant"] and content:
            messages.append({"role": role, "content": content})
            
    messages.append({"role": "user", "content": user_message})
    
    # === Provider routing ===
    provider    = data.get("provider", "local")       # local | apikey | memo
    req_api_key = data.get("api_key", "")              # only for apikey
    req_base_url= data.get("base_url", "")             # only for apikey
    req_model   = data.get("model", "")                # only for apikey

    if provider == "memo":
        # Route to Ulisse Memo cloud endpoint via OpenAI compatible client
        memo_url = f"{ULISSE_MEMO_URL}/v1"
        chat_client = OpenAI(
            base_url=memo_url,
            api_key=ULISSE_MEMO_BEARER or "memo-key"
        )
        chat_model = "deepseek-chat" # The proxy enforces the actual model

    elif provider == "apikey" and req_api_key:
        # User-supplied API key and base URL
        try:
            target_url = req_base_url or "https://api.openai.com/v1"
            if target_url.endswith('/'):
                target_url = target_url[:-1]

            user_client = OpenAI(
                api_key=req_api_key,
                base_url=target_url
            )
            chat_client = user_client
            chat_model  = req_model or "gpt-4o"
        except Exception as e:
            return jsonify({"error": f"Invalid API key config: {str(e)}"}), 400
    else:
        # Default: use .env / local model
        # Check if frontend sent local config overrides
        local_base_url = data.get("base_url")
        local_model    = data.get("model")
        
        if local_base_url:
            chat_client = OpenAI(
                api_key="ollama", # dummy for local
                base_url=local_base_url
            )
            chat_model = local_model or "llama3"
        else:
            chat_client = ai_client
            chat_model  = "deepseek-chat"

    # === LLM call (local / apikey paths) ===
    print(f"DEBUG: Provider={provider}, Model={chat_model}, BaseURL={chat_client.base_url}")
    print(f"DEBUG: Headers={chat_client.default_headers}")

    from flask import Response
    import copy
    
    def generate():
        nonlocal session_data, messages, sources
        
        # Salva la configurazione del modello nel contesto del thread corrente
        # così che i tool handler (come delegate_to_agno_agent) possano accedervi.
        agno_context.config = {
            "model": chat_model,
            "api_key": chat_client.api_key,
            "base_url": str(chat_client.base_url)
        }
        
        # Invia l'inizio della sessione
        yield f"data: {json.dumps({'event': 'session', 'session_id': session_id, 'session_title': session_data.get('title', 'Nuova Chat'), 'sources': sources})}\n\n"
        
        full_assistant_response = ""
        
        try:
            # Pass 1: Tool Routing
            yield f"data: {json.dumps({'event': 'think', 'message': 'Analisi in corso...'})}\n\n"
            
            # Check if we should enable reasoning for OpenRouter/DeepSeek-R1
            # Note: gpt-oss is V3 and doesn't natively support reasoning_content on all providers.
            extra_body = {}
            if "r1" in chat_model.lower():
                extra_body["reasoning"] = {"enabled": True}

            # Disable tools for OpenRouter :free models as they often don't support them
            current_tools = extra_tools if extra_tools else None
            is_openrouter = "openrouter.ai" in str(chat_client.base_url)
            
            if is_openrouter and ":free" in chat_model:
                current_tools = None
                print(f"DEBUG: Modello :free rilevato su OpenRouter, disabilitazione tools.")

            print(f"DEBUG: extra_body={extra_body}")

            # Use minimal parameters for OpenRouter to avoid routing conflicts
            if is_openrouter:
                response = chat_client.chat.completions.create(
                    model=chat_model,
                    messages=messages,
                    tools=current_tools,
                    extra_body=extra_body if extra_body else None
                )
            else:
                response = chat_client.chat.completions.create(
                    model=chat_model,
                    messages=messages,
                    temperature=0.7,
                    max_tokens=4000,
                    tools=current_tools,
                    extra_body=extra_body if extra_body else None
                )
            assistant_message = response.choices[0].message
            
            while assistant_message.tool_calls:
                # Convert to dict and ensure content is handled for API compatibility
                # Use exclude_none=True because Groq rejects explicit null values for fields like function_call
                msg_dict = assistant_message.model_dump(exclude_none=True)
                
                # Remove fields not supported in requests by some providers (like Groq)
                for field in ["annotations", "audio", "reasoning", "refusal"]:
                    msg_dict.pop(field, None)

                messages.append(msg_dict)

                for tool_call in assistant_message.tool_calls:
                    yield f"data: {json.dumps({'event': 'tool', 'message': f'Azione: {tool_call.function.name}'})}\n\n"
                    handler = extra_tool_handlers.get(tool_call.function.name)
                    
                    if handler:
                        try:
                            result = handler(tool_call.function.arguments)
                        except Exception as e:
                            result = f"Error executing tool {tool_call.function.name}: {str(e)}"
                    else:
                        result = f"Error: Tool '{tool_call.function.name}' not found."
                    
                    messages.append({
                        "role": "tool",
                        "tool_call_id": tool_call.id,
                        "name": tool_call.function.name,
                        "content": str(result)
                    })
                
                yield f"data: {json.dumps({'event': 'think', 'message': 'Elaborazione...'})}\n\n"
                if is_openrouter:
                    response = chat_client.chat.completions.create(
                        model=chat_model,
                        messages=messages,
                        tools=current_tools,
                        extra_body=extra_body if extra_body else None
                    )
                else:
                    response = chat_client.chat.completions.create(
                        model=chat_model,
                        messages=messages,
                        tools=current_tools,
                        extra_body=extra_body if extra_body else None
                    )
                assistant_message = response.choices[0].message
            
            # Pass 2: Generazione del testo finale (Stream)
            yield f"data: {json.dumps({'event': 'think', 'message': 'Rispondo...'})}\n\n"
            
            if is_openrouter:
                stream_response = chat_client.chat.completions.create(
                    model=chat_model,
                    messages=messages,
                    stream=True,
                    extra_body=extra_body if extra_body else None
                )
            else:
                stream_response = chat_client.chat.completions.create(
                    model=chat_model,
                    messages=messages,
                    temperature=0.7,
                    max_tokens=4000,
                    stream=True,
                    extra_body=extra_body if extra_body else None
                )
            
            for chunk in stream_response:
                if not chunk.choices:
                    continue
                delta = chunk.choices[0].delta
                
                # Check for reasoning_content (OpenRouter/DeepSeek-R1)
                reasoning = getattr(delta, 'reasoning_content', None) or (delta.model_dump().get('reasoning_content') if hasattr(delta, 'model_dump') else None)
                if reasoning:
                    yield f"data: {json.dumps({'event': 'think', 'message': reasoning, 'append': True})}\n\n"

                if delta.content:
                    full_assistant_response += delta.content
                    yield f"data: {json.dumps({'event': 'message', 'text': delta.content})}\n\n"
                    
        except Exception as e:
            print(f"LLM API error: {e}")
            yield f"data: {json.dumps({'event': 'error', 'message': f'Failed to generate response: {str(e)}'})}\n\n"
            if not full_assistant_response:
                full_assistant_response = "Errore di generazione."

        # --- Post-processing e Salvataggio ---
        timestamp = datetime.now().isoformat()
        
        # Auto-generate title from first message
        if len(chat_history) == 0:
            words = user_message.split()
            if len(words) >= 3:
                session_data["title"] = " ".join(w.capitalize() for w in words[:6])
            else:
                session_data["title"] = f"Chat {timestamp[:16]}"
        
        # AI Title refinement after 3 messages
        if len(chat_history) == 4: # 2 user + 2 assistant already, this is the 3rd pair
            try:
                title_prompt = [
                    {"role": "system", "content": "Generate a synthetic title (max 6 words) for this conversation between Ulisse and Toni. Respond ONLY with the title."},
                    {"role": "user", "content": f"Initial messages:\n" + "\n".join([f"{m['role']}: {m['content'][:100]}" for m in chat_history[:4]]) + f"\nuser: {user_message}"}
                ]
                t_resp = ai_client.chat.completions.create(
                    model="deepseek-chat",
                    messages=title_prompt,
                    max_tokens=20
                )
                new_title = t_resp.choices[0].message.content.strip().strip('"')
                if new_title:
                    session_data["title"] = new_title
            except Exception as e:
                print(f"Title generation error: {e}")

        # Append to session
        session_data["messages"].append({
            "role": "user",
            "content": user_message,
            "timestamp": timestamp
        })
        session_data["messages"].append({
            "role": "assistant",
            "content": full_assistant_response,
            "timestamp": timestamp,
            "sources": sources
        })
        session_data["updated_at"] = timestamp
        
        # Save session
        try:
            with open(get_session_path(session_id), "w", encoding="utf-8") as f:
                json.dump(session_data, f, ensure_ascii=False, indent=2)
        except Exception as e:
            print(f"Error saving session: {e}")

        # Save to legacy jsonl for RAG
        exchange = {
            "session_id": session_id,
            "title": session_data["title"],
            "timestamp": timestamp,
            "user_message": user_message,
            "assistant_response": full_assistant_response,
            "sources": sources
        }
        try:
            with open(new_conversations_file, "a", encoding="utf-8") as f:
                f.write(json.dumps(exchange, ensure_ascii=False) + "\n")
                
            # Real-time embedding into ChromaDB (Short-term memory)
            if chroma_status and collection is not None:
                document = f"User: {user_message}\nAssistant: {full_assistant_response}"
                chunk_id = f"chunk_{session_id}_{len(session_data['messages']) // 2}"
                collection.add(
                    documents=[document],
                    metadatas=[{
                        "session_id": session_id,
                        "title": session_data["title"],
                        "date": timestamp,
                        "type": "conversation"
                    }],
                    ids=[chunk_id]
                )
                print(f"Real-time STM update: Added chunk {chunk_id} to ChromaDB.")
                
        except Exception as e:
            print(f"Error saving legacy exchange or updating STM: {e}")

        # Segnale di fine stream
        yield f"data: {json.dumps({'event': 'done', 'session_title': session_data['title'], 'session_id': session_id})}\n\n"

    return Response(generate(), mimetype="text/event-stream")

@app.route("/sessions", methods=["GET"])
def get_sessions():
    sessions = []
    for f in sessions_dir.glob("*.json"):
        try:
            with open(f, "r", encoding="utf-8") as file:
                data = json.load(file)
                sessions.append({
                    "id": data["id"],
                    "title": data["title"],
                    "created_at": data.get("created_at"),
                    "updated_at": data.get("updated_at"),
                    "message_count": len(data.get("messages", [])) // 2
                })
        except:
            continue
    # Sort by updated_at descending
    sessions.sort(key=lambda x: x.get("updated_at", ""), reverse=True)
    return jsonify(sessions)

@app.route("/sessions/<session_id>", methods=["GET", "DELETE", "PATCH"])
def handle_session(session_id):
    path = get_session_path(session_id)
    if not path.exists():
        return jsonify({"error": "Session not found"}), 404
        
    if request.method == "GET":
        with open(path, "r", encoding="utf-8") as f:
            return jsonify(json.load(f))
            
    elif request.method == "DELETE":
        path.unlink()
        return jsonify({"status": "deleted"})
        
    elif request.method == "PATCH":
        data = request.json
        with open(path, "r", encoding="utf-8") as f:
            session_data = json.load(f)
        if "title" in data:
            session_data["title"] = data["title"]
        with open(path, "w", encoding="utf-8") as f:
            json.dump(session_data, f, ensure_ascii=False, indent=2)
        return jsonify({"status": "updated", "title": session_data["title"]})

@app.route("/memory/stats", methods=["GET"])
def get_memory_stats():
    sessions = []
    total_messages = 0
    all_titles = []
    
    for f in sessions_dir.glob("*.json"):
        try:
            with open(f, "r", encoding="utf-8") as file:
                data = json.load(file)
                sessions.append(data)
                total_messages += len(data.get("messages", []))
                all_titles.append(data.get("title", ""))
        except:
            continue
            
    if not sessions:
        return jsonify({
            "total_sessions": 0,
            "total_chunks": collection.count() if collection else 0,
            "total_messages": 0,
            "oldest_session": None,
            "newest_session": None,
            "most_active_topics": []
        })

    # Sort sessions by creation date
    sessions_sorted = sorted(sessions, key=lambda x: x.get("created_at", ""))
    
    # Extract keywords for topics
    words = []
    for title in all_titles:
        # Simple tokenization: lower case, alphanumeric only, length > 3
        found = re.findall(r'\b[a-z]{4,}\b', title.lower())
        words.extend(found)
    
    # Filter out common Italian stop words if necessary, but for now just top 10
    top_topics = [item[0] for item in Counter(words).most_common(10)]
    
    return jsonify({
        "total_sessions": len(sessions),
        "total_chunks": collection.count() if collection else 0,
        "total_messages": total_messages,
        "oldest_session": sessions_sorted[0].get("created_at"),
        "newest_session": sessions_sorted[-1].get("created_at"),
        "most_active_topics": top_topics
    })

@app.route("/memory/nodes", methods=["GET"])
def get_memory_nodes():
    nodes = []
    edges = []
    
    # Load Wiki Pages instead of sessions for the Memory View
    wiki_pages = list(wiki_pages_dir.glob("*.md"))
    
    for f in wiki_pages:
        try:
            title = f.stem.replace('_', ' ')
            content = f.read_text(encoding="utf-8")
            
            # Extract basic stats
            word_count = len(content.split())
            
            node = {
                "id": f.stem,
                "title": title,
                "content": content,
                "type": "wiki",
                "weight": max(5, min(word_count // 10, 20)),
                "created_at": datetime.fromtimestamp(f.stat().st_ctime).isoformat(),
                "updated_at": datetime.fromtimestamp(f.stat().st_mtime).isoformat()
            }
            nodes.append(node)
        except:
            continue
            
    # Generate Edges based on Wiki Links [[Page Name]]
    page_ids = [n["id"] for n in nodes]
    for n in nodes:
        content = n["content"]
        # Find all [[...]] links
        links = re.findall(r'\[\[([^\]]+)\]\]', content)
        for link in links:
            target_id = re.sub(r'[^\w\s-]', '', link).strip().replace(' ', '_')
            if target_id in page_ids and target_id != n["id"]:
                edges.append({
                    "source": n["id"],
                    "target": target_id,
                    "weight": 1
                })
                
    return jsonify({"nodes": nodes, "edges": edges})
    
@app.route("/memory/nodes/<node_id>", methods=["POST", "DELETE", "PATCH"])
def handle_memory_node(node_id):
    # node_id is the filename stem (e.g., 'Project_Alpha')
    file_path = wiki_pages_dir / f"{node_id}.md"
    
    if request.method == "DELETE":
        if file_path.exists():
            file_path.unlink()
            return jsonify({"status": "deleted"})
        return jsonify({"error": "Node not found"}), 404
        
    elif request.method == "POST":
        data = request.json
        new_content = data.get("content")
        if new_content is not None:
            file_path.write_text(new_content, encoding="utf-8")
            return jsonify({"status": "updated"})
        return jsonify({"error": "No content provided"}), 400

    elif request.method == "PATCH":
        data = request.json
        new_title = data.get("title")
        if new_title:
            # Generate new safe filename
            new_id = re.sub(r'[^\w\s-]', '', new_title).strip().replace(' ', '_')
            new_path = wiki_pages_dir / f"{new_id}.md"
            if file_path.exists():
                file_path.rename(new_path)
                return jsonify({"status": "renamed", "new_id": new_id})
            return jsonify({"error": "Node not found"}), 404
        return jsonify({"error": "No title provided"}), 400




@app.route("/sessions/<session_id>/related", methods=["GET"])
def get_related_sessions(session_id):
    path = get_session_path(session_id)
    if not path.exists() or not chroma_status or collection is None:
        return jsonify([])
        
    with open(path, "r", encoding="utf-8") as f:
        session_data = json.load(f)
    
    # Get last user message to find related
    msgs = session_data.get("messages", [])
    user_msgs = [m["content"] for m in msgs if m["role"] == "user"]
    if not user_msgs:
        return jsonify([])
    
    query = user_msgs[-1]
    try:
        results = collection.query(query_texts=[query], n_results=10)
        metas = results.get("metadatas", [[]])[0]
        
        related_ids = set()
        for m in metas:
            sid = m.get("session_id")
            if sid and sid != session_id:
                related_ids.add(sid)
        
        # Build session list
        related = []
        for sid in list(related_ids)[:5]:
            p = get_session_path(sid)
            if p.exists():
                with open(p, "r", encoding="utf-8") as f:
                    d = json.load(f)
                    related.append({"id": d["id"], "title": d["title"]})
        return jsonify(related)
    except:
        return jsonify([])

# Synthetic memory routes removed as per user request


@app.route("/graph", methods=["GET"])
def get_graph():
    try:
        sessions_list = []
        for f in sessions_dir.glob("*.json"):
            try:
                with open(f, "r", encoding="utf-8") as file:
                    sessions_list.append(json.load(file))
            except:
                continue
        
        nodes = []
        for s in sessions_list:
            sid = s["id"]
            title = s["title"]
            msgs = s.get("messages", [])
            msg_count = len(msgs) // 2
            
            nodes.append({
                "id": sid,
                "label": title,
                "type": "session",
                "weight": max(5, min(msg_count * 2, 20)),
                "created_at": s.get("created_at"),
                "updated_at": s.get("updated_at")
            })
            
        edges = []
        if chroma_status and collection is not None and len(sessions_list) > 1:
            # Generate edges based on RAG similarity
            for i, s1 in enumerate(sessions_list):
                sid1 = s1["id"]
                title1 = s1.get("title", "")
                
                try:
                    # Query ChromaDB for sessions similar to this one's title or first message
                    query_text = title1
                    if s1.get("messages"):
                        query_text += " " + s1["messages"][0]["content"]
                    
                    results = collection.query(
                        query_texts=[query_text],
                        n_results=5,
                        where={"session_id": {"$ne": sid1}}
                    )
                    
                    if results and results.get("metadatas"):
                        metas = results["metadatas"][0]
                        dists = results.get("distances", [[]])[0]
                        
                        for meta, dist in zip(metas, dists):
                            sid2 = meta.get("session_id")
                            if sid2:
                                # Connection strength based on distance (closer to 0 is stronger)
                                strength = max(0.1, 1.0 - dist)
                                edges.append({
                                    "from": sid1,
                                    "to": sid2,
                                    "strength": strength
                                })
                except Exception as e:
                    print(f"Error calculating graph edges for {sid1}: {e}")
        
        # Deduplicate edges (A->B and B->A might exist)
        unique_edges = {}
        for e in edges:
            pair = tuple(sorted([e["from"], e["to"]]))
            if pair not in unique_edges or e["strength"] > unique_edges[pair]["strength"]:
                unique_edges[pair] = e
        
        return jsonify({
            "nodes": nodes,
            "edges": list(unique_edges.values())
        })
    except Exception as e:
        print(f"Graph generation error: {e}")
        return jsonify({"error": str(e)}), 500

@app.route("/", methods=["GET"])
def index():
    return send_from_directory(str(frontend_dir), "index.html")

@app.route("/<path:filename>", methods=["GET"])
def serve_static(filename):
    return send_from_directory(str(frontend_dir), filename)

if __name__ == "__main__":
    # Disabilitiamo il reloader automatico per evitare conflitti con i processi in background (Playwright/Agno) su Windows
    app.run(host="0.0.0.0", port=5000, debug=True, use_reloader=False)
